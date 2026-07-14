# -*- coding: utf-8 -*-
"""
GlamHub Presentation Generator  — v3  (20 slides, student-grade)
Reference style: GLAMHUB-Modern-Cosmetics-E-Commerce (1).pptx
Spec: Heading 36-40pt bold  |  Body 28pt  |  10" × 5.625"
"""

import io, os, urllib.request
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─────────────────────────── SLIDE SIZE ────────────────────────────────────
W, H = 10.0, 5.625   # widescreen 16:9

# ─────────────────────────── COLOUR PALETTE ────────────────────────────────
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK   = RGBColor(0x00, 0x00, 0x00)
C_DARK    = RGBColor(0x27, 0x25, 0x25)   # body text
C_BLUE    = RGBColor(0x49, 0x50, 0xBC)   # section tag
C_LAVEND  = RGBColor(0xDA, 0xDB, 0xF1)   # card / pill
C_CARD2   = RGBColor(0xFF, 0xF8, 0xFA)   # soft pink card
C_ACCENT  = RGBColor(0xE0, 0x5A, 0x88)   # rose pink
C_PRIM    = RGBColor(0x8A, 0x1C, 0x4A)   # deep burgundy
C_BGLIGHT = RGBColor(0xF9, 0xF9, 0xFF)   # very light bg

# ─────────────────── HIGH-QUALITY, SLIDE-RELEVANT IMAGES ───────────────────
# Every URL targets a specific cosmetics/tech concept and is freely usable
IMGS = {
    # Slide 1 – Title: glamorous beauty flat-lay
    "title":   "https://images.unsplash.com/photo-1522335789203-aabd1fc54bc9?w=900&q=85",
    # Slide 2 – Introduction: woman shopping online with cosmetics
    "intro":   "https://images.unsplash.com/photo-1607748851687-ba9a10438621?w=900&q=85",
    # Slide 3 – Problem: frustrated customer, broken experience
    "problem": "https://images.unsplash.com/photo-1556228720-195a672e8a03?w=900&q=85",
    # Slide 4 – Objectives: goal/target concept
    "obj":     "https://images.unsplash.com/photo-1512496015851-a90fb38ba796?w=900&q=85",
    # Slide 5 – Target audience: diverse beauty consumers
    "target":  "https://images.unsplash.com/photo-1487412947147-5cebf100ffc2?w=900&q=85",
    # Slide 6 – Architecture: clean code / tech stack
    "arch":    "https://images.unsplash.com/photo-1461749280684-dccba630e2f6?w=900&q=85",
    # Slide 7 – Methodology: agile / sprint planning whiteboard
    "dev":     "https://images.unsplash.com/photo-1454165804606-c3d57bc86b40?w=900&q=85",
    # Slide 8 – Authentication: login / password security
    "auth":    "https://images.unsplash.com/photo-1614680376573-df3480f0c6ff?w=900&q=85",
    # Slide 9A – Catalog grid: skincare products arranged beautifully
    "cat_a":   "https://images.unsplash.com/photo-1596462502278-27bfdc403348?w=900&q=85",
    # Slide 9B – Catalog grid: makeup brushes and palette
    "cat_b":   "https://images.unsplash.com/photo-1512496015851-a90fb38ba796?w=900&q=85",
    # Slide 9C – Catalog grid: hair care and body products
    "cat_c":   "https://images.unsplash.com/photo-1571781926291-c477ebfd024b?w=900&q=85",
    # Slide 10 – Cart: shopping bags / cart concept
    "cart":    "https://images.unsplash.com/photo-1607082349566-187342175e2f?w=900&q=85",
    # Slide 11 – Wishlist: heart/love concept with cosmetics
    "wish":    "https://images.unsplash.com/photo-1556228578-8c89e6adf883?w=900&q=85",
    # Slide 12 – Payment: contactless / card payment
    "pay":     "https://images.unsplash.com/photo-1563013544-824ae1b704d3?w=900&q=85",
    # Slide 13 – Orders: parcel / delivery concept
    "order":   "https://images.unsplash.com/photo-1530021232320-687d8e3dba54?w=900&q=85",
    # Slide 14 – Database: data visualization / ER concept
    "db":      "https://images.unsplash.com/photo-1544383835-bda2bc66a55d?w=900&q=85",
    # Slide 15 – UX: mobile responsive design in hand
    "ux":      "https://images.unsplash.com/photo-1555421689-3f034debb7a6?w=900&q=85",
    # Slide 16 – Security: padlock / cyber security
    "sec":     "https://images.unsplash.com/photo-1563986768494-4dee2763ff3f?w=900&q=85",
    # Slide 17 – Performance: speed / analytics dashboard
    "perf":    "https://images.unsplash.com/photo-1551288049-bebda4e38f71?w=900&q=85",
    # Slide 18 – Validation: checklist / testing
    "val":     "https://images.unsplash.com/photo-1460925895917-afdab827c52f?w=900&q=85",
    # Slide 19 – Business: revenue / growth chart
    "biz":     "https://images.unsplash.com/photo-1605810230434-7631ac76ec81?w=900&q=85",
    # Slide 20 – Conclusion: beauty + future (sunrise / glow)
    "concl":   "https://images.unsplash.com/photo-1487412947147-5cebf100ffc2?w=900&q=85",
}

# ═══════════════════════════ HELPER UTILITIES ═══════════════════════════════

def fetch_image(url):
    """Download image from URL → BytesIO; returns None on failure."""
    try:
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=12) as r:
            return io.BytesIO(r.read())
    except Exception as e:
        print(f"  [WARN] Could not fetch {url}: {e}")
        return None


def add_picture(slide, url, l, t, w, h):
    """Place an image on the slide (silently skips on download failure)."""
    img = fetch_image(url)
    if img:
        slide.shapes.add_picture(img, Inches(l), Inches(t), Inches(w), Inches(h))


def add_rect(slide, l, t, w, h, fill, line=None, lw=None, rounded=False):
    """Add a filled rectangle (or rounded rect) shape."""
    shape_id = 5 if rounded else 1
    s = slide.shapes.add_shape(shape_id, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        if lw:
            s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s


def tb(slide, l, t, w, h, text, fname, fsize, color,
       bold=False, align=PP_ALIGN.LEFT, wrap=True):
    """Add a single-paragraph textbox and return (shape, text_frame)."""
    tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.text = text
    r.font.name = fname
    r.font.size = Pt(fsize)
    r.font.bold = bold
    r.font.color.rgb = color
    return tx, tf


def section_pill(slide, label, l, t):
    """Small lavender pill badge with section label text."""
    add_rect(slide, l, t, 1.60, 0.22, C_LAVEND, rounded=True)
    tb(slide, l + 0.07, t + 0.01, 1.50, 0.20, label, "Inter", 9, C_DARK)


def heading(slide, text, l, t, w, size=38):
    """Bold slide heading at specified position."""
    tb(slide, l, t, w, 0.62, text, "Inter Bold", size, C_BLACK, bold=True)


def label_detail_block(slide, l, t, w, h, items, label_size=14, detail_size=11):
    """
    Render a stacked list of (label, detail) pairs inside a textbox.
    label_size / detail_size in points.
    """
    tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    for i, (label, detail) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        r = p.runs[0] if p.runs else p.add_run()
        r.text = label
        r.font.name = "Inter Bold"
        r.font.size = Pt(label_size)
        r.font.bold = True
        r.font.color.rgb = C_BLACK
        if detail:
            p2 = tf.add_paragraph()
            p2.space_after = Pt(9)
            r2 = p2.runs[0] if p2.runs else p2.add_run()
            r2.text = detail
            r2.font.name = "Inter"
            r2.font.size = Pt(detail_size)
            r2.font.bold = False
            r2.font.color.rgb = C_DARK
    return tf


def bullet_box(slide, l, t, w, h, items, size=11):
    """Render a bullet list (•  text) inside a textbox."""
    tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        r = p.runs[0] if p.runs else p.add_run()
        r.text = f"\u2022  {text}"
        r.font.name = "Inter"
        r.font.size = Pt(size)
        r.font.bold = False
        r.font.color.rgb = C_DARK
    return tf


def three_cards(slide, cards, top=2.5, card_w=2.9, card_h=1.8,
                start_l=0.54, gap=0.19, body_size=11):
    """3 equal lavender cards side by side. cards = [(title, body), ...]"""
    for i, (title, body) in enumerate(cards):
        l = start_l + i * (card_w + gap)
        add_rect(slide, l, top, card_w, card_h, C_LAVEND)
        tb(slide, l + 0.15, top + 0.15, card_w - 0.3, 0.28,
           title, "Inter Bold", 14, C_DARK, bold=True)
        tb(slide, l + 0.15, top + 0.52, card_w - 0.3, card_h - 0.65,
           body, "Inter", body_size, C_DARK, wrap=True)


def right_image_slide(prs, blank, img_key, section, title_text, items, slide_num):
    """
    Standard template: image on right (3.75"), text on left.
    items = [(label, detail), ...]  — uses label_detail_block.
    """
    print(f"Building Slide {slide_num}: {title_text}...")
    slide = prs.slides.add_slide(blank)
    add_picture(slide, IMGS[img_key], 6.25, 0.0, 3.75, 5.625)
    section_pill(slide, section, 0.54, 0.45)
    heading(slide, title_text, 0.54, 0.73, 5.50)
    label_detail_block(slide, 0.54, 1.50, 5.30, 3.80, items, 14, 11)
    return slide


def left_image_slide(prs, blank, img_key, section, title_text, items, slide_num):
    """
    Standard template: image on left (3.75"), text on right.
    items = [(label, detail), ...]  — uses label_detail_block.
    """
    print(f"Building Slide {slide_num}: {title_text}...")
    slide = prs.slides.add_slide(blank)
    add_picture(slide, IMGS[img_key], 0.0, 0.0, 3.75, 5.625)
    section_pill(slide, section, 4.29, 0.45)
    heading(slide, title_text, 4.29, 0.73, 5.30)
    label_detail_block(slide, 4.29, 1.50, 5.30, 3.80, items, 14, 11)
    return slide


# ═══════════════════════════ BUILD PRESENTATION ═════════════════════════════

def create_presentation():
    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    # ──────────────────────────────────────────────────────────────────────
    # SLIDE 1 — TITLE
    # ──────────────────────────────────────────────────────────────────────
    print("Building Slide 1: Title...")
    slide = prs.slides.add_slide(blank)

    add_picture(slide, IMGS["title"], 6.25, 0.0, 3.75, 5.625)
    section_pill(slide, "COMPUTER SCIENCE & ENGINEERING", 0.54, 1.55)

    tx, tf = tb(slide, 0.54, 1.90, 5.16, 1.20,
                "GLAMHUB: Modern Cosmetics\nE-Commerce",
                "Inter Bold", 38, C_BLACK, bold=True)
    tf.word_wrap = True

    tx2 = slide.shapes.add_textbox(Inches(0.54), Inches(3.25), Inches(5.16), Inches(1.00))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    info_lines = [
        "Project Supervisor:  [Insert Supervisor Name]",
        "Developed By:  [Insert Student Names]",
        "Department of Computer Science & Engineering",
    ]
    for i, line in enumerate(info_lines):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        r = p.runs[0] if p.runs else p.add_run()
        r.text = line
        r.font.name = "Inter"
        r.font.size = Pt(13)
        r.font.color.rgb = C_DARK

    # ──────────────────────────────────────────────────────────────────────
    # SLIDE 2 — Introduction to GlamHub
    # ──────────────────────────────────────────────────────────────────────
    print("Building Slide 2: Introduction...")
    slide = prs.slides.add_slide(blank)

    section_pill(slide, "INTRODUCTION", 0.54, 1.10)
    add_rect(slide, 2.00, 1.10, 1.65, 0.22, C_LAVEND, rounded=True)
    tb(slide, 2.07, 1.11, 1.60, 0.20, "PROJECT SUMMARY", "Inter", 9, C_BLUE)
    heading(slide, "Introduction to GlamHub", 0.54, 1.38, 4.80)

    tx = slide.shapes.add_textbox(Inches(0.54), Inches(2.08), Inches(9.0), Inches(0.52))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(8)
    r = p.runs[0] if p.runs else p.add_run()
    r.text = ("A specialized e-commerce web application purpose-built for beauty and cosmetics retail, "
              "GlamHub combines a robust Django backend with a modern, responsive front-end.")
    r.font.name = "Inter"
    r.font.size = Pt(11)
    r.font.color.rgb = C_DARK

    three_cards(slide, [
        ("Tech Stack",
         "Python, Django MVT, HTML5, CSS3, JavaScript — production-ready from day one."),
        ("User Experience",
         "Seamless, responsive, and highly secure shopping experience across all devices."),
        ("Domain Focus",
         "Cosmetics-first design: rich product metadata, labels, and intelligent search."),
    ], top=2.78, card_w=2.90, card_h=1.75)

    # ──────────────────────────────────────────────────────────────────────
    # SLIDE 3 — Problem Statement  (image right)
    # ──────────────────────────────────────────────────────────────────────
    print("Building Slide 3: Problem Statement...")
    slide = prs.slides.add_slide(blank)
    add_picture(slide, IMGS["problem"], 6.25, 0.0, 3.75, 5.625)
    section_pill(slide, "INTRODUCTION", 0.54, 0.52)
    heading(slide, "Problem Statement", 0.54, 0.83, 4.50)

    problems = [
        ("Generic Experience",
         "General e-commerce platforms lack specialized UX for cosmetic shoppers who need detailed product info."),
        ("Cart Abandonment",
         "High abandonment rates persist due to complex checkout flows and insufficient engagement features."),
        ("Trust & Security",
         "Consumers hesitate to purchase online without secure payment systems and clear data protection."),
    ]
    col_w = 1.80
    gap   = 0.17
    top   = 1.55
    for i, (title, body) in enumerate(problems):
        l = 0.54 + i * (col_w + gap)
        tb(slide, l, top, col_w, 0.35, title, "Inter Bold", 13, C_BLACK, bold=True)
        tb(slide, l, top + 0.42, col_w, 2.60, body, "Inter", 10, C_DARK, wrap=True)

    # ──────────────────────────────────────────────────────────────────────
    # SLIDE 4 — Project Objectives  (image right)
    # ──────────────────────────────────────────────────────────────────────
    print("Building Slide 4: Project Objectives...")
    slide = prs.slides.add_slide(blank)
    add_picture(slide, IMGS["obj"], 6.25, 0.0, 3.75, 5.625)
    section_pill(slide, "INTRODUCTION", 0.54, 0.52)
    heading(slide, "Project Objectives", 0.54, 0.83, 5.0)
    label_detail_block(slide, 0.54, 1.55, 5.30, 3.80, [
        ("Scalable Backend",  "Build a robust, scalable backend using Django's MVT architecture."),
        ("Intuitive UI",      "Design an appealing, user-friendly interface tailored for cosmetics browsing."),
        ("Secure Payments",   "Integrate Razorpay for secure, reliable INR payment processing."),
        ("Personalization",   "Implement wishlists, order history, and customer profile management."),
    ])

    # ──────────────────────────────────────────────────────────────────────
    # SLIDE 5 — Target Market & Audience  (image right)
    # ──────────────────────────────────────────────────────────────────────
    print("Building Slide 5: Target Audience...")
    slide = prs.slides.add_slide(blank)
    add_picture(slide, IMGS["target"], 6.25, 0.0, 3.75, 5.625)
    section_pill(slide, "INTRODUCTION", 0.54, 0.73)
    heading(slide, "Target Market & Audience", 0.54, 1.01, 5.50)

    tb(slide, 0.54, 1.78, 2.68, 0.28, "Who We Serve", "Inter Bold", 13, C_BLACK, bold=True)
    tb(slide, 0.54, 2.14, 2.68, 2.20,
       ("Skincare, makeup, and hair care enthusiasts who demand detailed "
        "product information and intuitive discovery.\n\n"
        "Beauty consumers who value curated, premium product selections."),
       "Inter", 10, C_DARK, wrap=True)

    add_rect(slide, 3.40, 1.66, 2.72, 3.24, C_CARD2)
    tb(slide, 3.56, 1.82, 2.44, 0.28, "Core Goals", "Inter Bold", 13, C_BLACK, bold=True)
    bullet_box(slide, 3.56, 2.18, 2.44, 2.50, [
        "Bridge the gap between customers and premium cosmetic brands.",
        "Provide rich product metadata: labels, categories, search.",
        "Create a personalised discovery experience that drives engagement.",
    ], size=10)

    # ──────────────────────────────────────────────────────────────────────
    # SLIDE 6 — System Architecture
    # ──────────────────────────────────────────────────────────────────────
    print("Building Slide 6: System Architecture...")
    slide = prs.slides.add_slide(blank)
    section_pill(slide, "DESIGN", 0.54, 0.37)
    heading(slide, "System Architecture", 0.54, 0.62, 4.50)
    add_picture(slide, IMGS["arch"], 5.10, 0.60, 4.36, 4.36)
    label_detail_block(slide, 0.54, 1.35, 4.36, 3.80, [
        ("MVT Pattern",  "Django's Model-View-Template separates data, business logic, and presentation cleanly."),
        ("Database",     "SQLite3 with normalized schema and FK relationships for data integrity."),
        ("Security",     "CSRF & XSS protection via Django middleware; PBKDF2 password hashing built in."),
        ("API Layer",    "Razorpay API integrated via secure server-side calls with JSON response handling."),
    ])

    # ──────────────────────────────────────────────────────────────────────
    # SLIDE 7 — Development Methodology  (image left)
    # ──────────────────────────────────────────────────────────────────────
    print("Building Slide 7: Development Methodology...")
    slide = prs.slides.add_slide(blank)
    add_picture(slide, IMGS["dev"], 0.0, 0.0, 3.75, 5.625)
    section_pill(slide, "DESIGN", 4.29, 0.40)
    heading(slide, "Development Methodology", 4.29, 0.68, 5.0)

    tx = slide.shapes.add_textbox(Inches(4.29), Inches(1.38), Inches(5.16), Inches(0.50))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.runs[0] if p.runs else p.add_run()
    r.text = ("GlamHub was developed using an Agile methodology with iterative sprint cycles, "
              "ensuring continuous integration of features and user feedback.")
    r.font.name = "Inter"
    r.font.size = Pt(10)
    r.font.color.rgb = C_DARK

    label_detail_block(slide, 4.29, 2.00, 5.30, 3.40, [
        ("Sequential Design",   "DB design → view logic → template rendering — a clean, logical flow."),
        ("Iterative Sprints",   "Features developed in focused sprint cycles for continuous improvement."),
        ("Systematic Testing",  "User workflows and payment integration undergo rigorous validation cycles."),
    ])

    # ──────────────────────────────────────────────────────────────────────
    # SLIDE 8 — Customer Authentication
    # ──────────────────────────────────────────────────────────────────────
    print("Building Slide 8: Customer Authentication...")
    slide = prs.slides.add_slide(blank)
    section_pill(slide, "KEY FEATURES", 0.54, 0.99)
    heading(slide, "Customer Authentication", 0.54, 1.26, 4.50)
    add_picture(slide, IMGS["auth"], 0.54, 2.08, 3.80, 2.20)

    tx = slide.shapes.add_textbox(Inches(4.73), Inches(1.95), Inches(4.80), Inches(0.70))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.space_after = Pt(8)
    r = p.runs[0] if p.runs else p.add_run()
    r.text = ("GlamHub implements a custom Customer model extending Django's "
              "AbstractUser, enabling mobile number storage alongside standard auth fields.")
    r.font.name = "Inter"
    r.font.size = Pt(10)
    r.font.color.rgb = C_DARK

    bullet_box(slide, 4.73, 2.82, 4.80, 2.50, [
        "Secure login, registration, and logout with form validation.",
        "Session-based authentication maintains state across navigations.",
        "Passwords stored using PBKDF2 hashing — industry standard.",
        "login_required decorators protect all user-sensitive operations.",
    ], size=10)

    # ──────────────────────────────────────────────────────────────────────
    # SLIDE 9 — Product Catalog & Discovery  (3-image grid)
    # ──────────────────────────────────────────────────────────────────────
    print("Building Slide 9: Product Catalog...")
    slide = prs.slides.add_slide(blank)
    section_pill(slide, "KEY FEATURES", 0.54, 0.80)
    heading(slide, "Product Catalog & Discovery", 0.54, 1.08, 5.50)

    cat_imgs = [IMGS["cat_a"], IMGS["cat_b"], IMGS["cat_c"]]
    cat_items = [
        ("Dynamic Categories",
         "Skincare, Makeup, Haircare, Fragrances, and Bodycare — each with curated subcategories."),
        ("Product Labels",
         "Best Seller, New Arrival, Price Drop, and BOGO tags for rapid discovery."),
        ("Real-Time Search",
         "Dynamic search across the entire catalog returns instant, relevant results."),
    ]
    img_w = 2.84
    img_h = 1.78
    gap   = 0.19
    for i, ((title, body), url) in enumerate(zip(cat_items, cat_imgs)):
        l = 0.54 + i * (img_w + gap)
        add_picture(slide, url, l, 1.72, img_w, img_h)
        tb(slide, l, 3.62, img_w - 0.10, 0.28, title, "Inter Bold", 13, C_DARK, bold=True)
        tb(slide, l, 3.96, img_w,         0.90, body,  "Inter",     10, C_DARK, wrap=True)

    # ──────────────────────────────────────────────────────────────────────
    # SLIDE 10 — Interactive Shopping Cart  (image left)
    # ──────────────────────────────────────────────────────────────────────
    print("Building Slide 10: Shopping Cart...")
    slide = prs.slides.add_slide(blank)
    add_picture(slide, IMGS["cart"], 0.0, 0.0, 3.75, 5.625)
    section_pill(slide, "KEY FEATURES", 4.29, 0.45)
    heading(slide, "Interactive Shopping Cart", 4.29, 0.73, 5.30)

    # Intro paragraph with bold inline
    tx = slide.shapes.add_textbox(Inches(4.29), Inches(1.45), Inches(2.45), Inches(1.10))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.space_after = Pt(6)
    for txt_part, is_bold in [
        ("GlamHub manages the cart via a dedicated ", False),
        ("CartItem model", True),
        (", persisting state in the database for a reliable experience.\n\n"
         "Users can modify quantities or remove items — no page reload required.", False),
    ]:
        r = p.add_run()
        r.text = txt_part
        r.font.name = "Inter Bold" if is_bold else "Inter"
        r.font.size = Pt(10)
        r.font.bold = is_bold
        r.font.color.rgb = C_DARK

    # Capabilities card
    add_rect(slide, 6.96, 1.32, 2.65, 3.96, C_CARD2)
    tb(slide, 7.12, 1.47, 2.44, 0.28, "Cart Capabilities", "Inter Bold", 13, C_BLACK, bold=True)
    bullet_box(slide, 7.12, 1.90, 2.44, 3.20, [
        "Database-driven cart via CartItem model.",
        "Dynamic quantity updates & item removal.",
        "Automatic subtotal calculations.",
        "70% Pink Summer Sale discount applied.",
        "Seamless integration with checkout flow.",
    ], size=10)

    # ──────────────────────────────────────────────────────────────────────
    # SLIDE 11 — Wishlist & Personalization  (image right)
    # ──────────────────────────────────────────────────────────────────────
    right_image_slide(prs, blank, "wish", "KEY FEATURES",
                      "Wishlist & Personalization", [
        ("Wishlist Model",      "Dedicated WishlistItem model saves favourite products per user."),
        ("One-Click Toggle",    "Add or remove items from any page with instant visual feedback."),
        ("Account Dashboard",   "Customer profile page allows updating personal details securely."),
        ("Password Management", "Secure password-change form with session re-authentication."),
    ], 11)

    # ──────────────────────────────────────────────────────────────────────
    # SLIDE 12 — Payment Gateway Integration  (image left)
    # ──────────────────────────────────────────────────────────────────────
    left_image_slide(prs, blank, "pay", "KEY FEATURES",
                     "Payment Gateway Integration", [
        ("Razorpay API",       "Full Razorpay API integration for seamless INR payment processing."),
        ("Server-Side Orders", "Secure order object created server-side before payment initiation."),
        ("Verification",       "Payment callback verifies transaction before creating Order records."),
        ("PCI Compliance",     "No raw card data stored — all handled by Razorpay's secure servers."),
    ], 12)

    # ──────────────────────────────────────────────────────────────────────
    # SLIDE 13 — Order Management & Tracking  (image right)
    # ──────────────────────────────────────────────────────────────────────
    right_image_slide(prs, blank, "order", "KEY FEATURES",
                      "Order Management & Tracking", [
        ("Order Model",     "Persistent Order model stores product, quantity, address, payment ID, status."),
        ("Auto-Generation", "Orders auto-created and cart auto-cleared upon successful payment."),
        ("Order History",   "User-accessible My Orders page shows complete purchase history."),
        ("Payment Status",  "Tracks payment_status and payment_id per order for reconciliation."),
    ], 13)

    # ──────────────────────────────────────────────────────────────────────
    # SLIDE 14 — Database Design & Integrity  (image right + 3 cards)
    # ──────────────────────────────────────────────────────────────────────
    print("Building Slide 14: Database Design...")
    slide = prs.slides.add_slide(blank)
    section_pill(slide, "RESULTS", 0.54, 0.45)
    heading(slide, "Database Design & Integrity", 0.54, 0.73, 5.50)
    add_picture(slide, IMGS["db"], 6.25, 0.0, 3.75, 5.625)
    three_cards(slide, [
        ("Normalised Schema",
         "3NF normalization eliminates data redundancy and update/insert anomalies across all models."),
        ("Referential Integrity",
         "FK constraints ensure consistent relationships: Customer ↔ Item ↔ Order ↔ CartItem."),
        ("Query Efficiency",
         "Optimized ORM queries with indexed fields reduce database response time under load."),
    ], top=1.68, card_w=2.90, card_h=2.20)

    # ──────────────────────────────────────────────────────────────────────
    # SLIDE 15 — User Experience & Responsiveness  (image left)
    # ──────────────────────────────────────────────────────────────────────
    left_image_slide(prs, blank, "ux", "RESULTS",
                     "User Experience & Responsiveness", [
        ("Responsive Design",    "Fully responsive layout adapts to mobile, tablet, and desktop screens."),
        ("Intuitive Navigation", "Clean category structure reduces cognitive load and aids discoverability."),
        ("Instant Feedback",     "Django flash messaging provides success/error alerts after every action."),
        ("Dynamic Interactions", "Cart and wishlist interactions feel fluid with zero page reloads."),
    ], 15)

    # ──────────────────────────────────────────────────────────────────────
    # SLIDE 16 — Security & Reliability Outcomes  (image right)
    # ──────────────────────────────────────────────────────────────────────
    right_image_slide(prs, blank, "sec", "RESULTS",
                      "Security & Reliability Outcomes", [
        ("CSRF Protection",   "All POST routes protected by Django's CSRF middleware — zero vulnerabilities."),
        ("Password Security", "Passwords stored as PBKDF2 hashes with salting; never in plain text."),
        ("Payment Integrity", "No order is created unless Razorpay payment verification succeeds."),
        ("Access Control",    "login_required decorators prevent unauthorized access to all user views."),
    ], 16)

    # ──────────────────────────────────────────────────────────────────────
    # SLIDE 17 — Performance Benchmarks  (image left)
    # ──────────────────────────────────────────────────────────────────────
    left_image_slide(prs, blank, "perf", "RESULTS",
                     "Performance Benchmarks", [
        ("Optimised Queries", "Efficient Django ORM queries minimise database round-trips per request."),
        ("Low Latency",       "Minimal response time for cart and wishlist CRUD operations."),
        ("Media Handling",    "Image delivery via pre-defined upload paths — CDN-ready for scale."),
        ("Search Speed",      "Indexed name field ensures rapid icontains search across the full catalog."),
    ], 17)

    # ──────────────────────────────────────────────────────────────────────
    # SLIDE 18 — Functional Validation Results  (image right + 3 cards)
    # ──────────────────────────────────────────────────────────────────────
    print("Building Slide 18: Functional Validation...")
    slide = prs.slides.add_slide(blank)
    section_pill(slide, "RESULTS", 0.54, 0.45)
    heading(slide, "Functional Validation Results", 0.54, 0.73, 6.0)
    add_picture(slide, IMGS["val"], 6.25, 0.0, 3.75, 5.625)
    three_cards(slide, [
        ("End-to-End Flow",
         "Validated: Register → Login → Browse → Add to Cart → Checkout → Pay → Orders."),
        ("Unit Testing",
         "Django test suite confirms model constraints, form validation, and URL routing."),
        ("Edge Cases",
         "Handled: empty-cart checkout, duplicate wishlist entries, varied price formats."),
    ], top=1.68, card_w=2.90, card_h=2.20)

    # ──────────────────────────────────────────────────────────────────────
    # SLIDE 19 — Business & Operational Value  (image left)
    # ──────────────────────────────────────────────────────────────────────
    left_image_slide(prs, blank, "biz", "RESULTS",
                     "Business & Operational Value", [
        ("Sales Automation", "Automates the full sales cycle — reduces manual operational overhead."),
        ("Intent Capture",   "Wishlist data captures purchase intent for targeted marketing campaigns."),
        ("Discount Engine",  "Promotional discount system (Summer Sale) actively drives conversions."),
        ("Scalability",      "Architecture is ready for analytics dashboards and CRM integrations."),
    ], 19)

    # ──────────────────────────────────────────────────────────────────────
    # SLIDE 20 — Conclusion & Future Scope  (image right)
    # ──────────────────────────────────────────────────────────────────────
    print("Building Slide 20: Conclusion...")
    slide = prs.slides.add_slide(blank)
    add_picture(slide, IMGS["concl"], 6.25, 0.0, 3.75, 5.625)
    section_pill(slide, "CONCLUSION", 0.54, 0.52)
    heading(slide, "Conclusion & Future Scope", 0.54, 0.80, 5.50)

    tx = slide.shapes.add_textbox(Inches(0.54), Inches(1.52), Inches(5.30), Inches(0.52))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.space_after = Pt(8)
    r = p.runs[0] if p.runs else p.add_run()
    r.text = ("GlamHub is a fully functional, secure, and scalable e-commerce platform that "
              "successfully bridges the gap in beauty retail online.")
    r.font.name = "Inter"
    r.font.size = Pt(10)
    r.font.color.rgb = C_DARK

    three_cards(slide, [
        ("Achieved",
         "Secure auth, smart catalog, cart, wishlist, Razorpay payments, and order tracking — all delivered."),
        ("Engineering Value",
         "Real-world Django + API integration with clean MVT architecture and security-first design."),
        ("Future Scope",
         "AI recommendations, multi-currency support, logistics API, seller dashboard & advanced analytics."),
    ], top=2.25, card_w=2.90, card_h=2.35, start_l=0.54)

    # ──────────────────────────────────────────────────────────────────────
    # SAVE
    # ──────────────────────────────────────────────────────────────────────
    output = r"d:\GlamHub\GlamHub_Final.pptx"
    prs.save(output)
    print(f"\nDONE: Saved -> {output}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    create_presentation()
