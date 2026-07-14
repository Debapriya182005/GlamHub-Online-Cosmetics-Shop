from pptx import Presentation
from pptx.util import Inches, Pt

REF = r'd:\GlamHub\GLAMHUB-Modern-Cosmetics-E-Commerce (1).pptx'
prs = Presentation(REF)

print(f'Slide size: {prs.slide_width.inches:.3f} x {prs.slide_height.inches:.3f} inches')
print(f'Total slides: {len(prs.slides)}')

# Read slides 1-2 (title) and 14-20 (results + conclusion)
target = list(range(0, 3)) + list(range(13, len(prs.slides)))

for i in target:
    if i >= len(prs.slides):
        break
    slide = prs.slides[i]
    print(f'\n{"="*70}')
    print(f'SLIDE {i+1}')
    print(f'{"="*70}')

    bg = slide.background.fill
    try:
        print(f'  Background: #{bg.fore_color.rgb}')
    except Exception:
        try:
            print(f'  Background type: {bg.type}')
        except Exception:
            print('  Background: unknown')

    for shape in slide.shapes:
        l = shape.left / 914400 if shape.left else 0
        t = shape.top / 914400 if shape.top else 0
        w = shape.width / 914400 if shape.width else 0
        h = shape.height / 914400 if shape.height else 0

        fill_info = ''
        try:
            if shape.fill.type == 1:
                fill_info = f' fill=#{shape.fill.fore_color.rgb}'
            elif shape.fill.type not in (None, 5):
                fill_info = f' fill_type={shape.fill.type}'
        except Exception:
            pass

        print(f'  [Shape] "{shape.name}" type={shape.shape_type} '
              f'pos=({l:.3f}",{t:.3f}") size=({w:.3f}"x{h:.3f}"){fill_info}')

        if shape.has_text_frame:
            for pi, para in enumerate(shape.text_frame.paragraphs):
                txt = para.text.strip()
                if not txt:
                    continue
                print(f'    Para[{pi}]: "{txt[:100]}"')
                for ri, run in enumerate(para.runs):
                    rf = run.font
                    try:
                        rc = str(rf.color.rgb)
                    except Exception:
                        rc = 'inherited'
                    print(f'      Run[{ri}]: name="{rf.name}" size={rf.size} '
                          f'bold={rf.bold} color={rc}')
